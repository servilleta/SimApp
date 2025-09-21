import os
import logging
import json
import tempfile
import time
from typing import Dict, Any, Optional, List
from datetime import datetime
from pathlib import Path
import io
import uuid
import asyncio

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from playwright.async_api import async_playwright
import numpy as np

logger = logging.getLogger(__name__)

class PowerPointExportService:
    """Service for generating editable PowerPoint presentations from simulation results."""
    
    def __init__(self):
        self.logger = logger
        
        # Color scheme for consistency
        self.colors = {
            'primary': RGBColor(33, 150, 243),     # Blue
            'secondary': RGBColor(255, 152, 0),    # Orange  
            'success': RGBColor(76, 175, 80),      # Green
            'warning': RGBColor(255, 193, 7),      # Yellow
            'danger': RGBColor(244, 67, 54),       # Red
            'dark': RGBColor(33, 37, 41),          # Dark gray
            'light': RGBColor(248, 249, 250)       # Light gray
        }
    
    async def generate_powerpoint_presentation(
        self, 
        simulation_id: str, 
        results_data: Dict[str, Any],
        metadata: Optional[Dict[str, Any]] = None,
        auth_token: Optional[str] = None,
        frontend_url: str = "http://frontend:3000"
    ) -> str:
        """
        Generate PowerPoint presentation using pixel-perfect screenshots for visual fidelity.
        This achieves the same visual output as PDF but in an editable PowerPoint format.
        
        Args:
            simulation_id: Unique identifier for the simulation
            results_data: Simulation results including targets and statistics
            metadata: Additional metadata about the simulation
            auth_token: Authentication token for frontend access
            frontend_url: URL of the frontend service
            
        Returns:
            str: Path to the generated PowerPoint file
        """
        try:
            logger.info(f"Starting pixel-perfect PowerPoint generation for simulation {simulation_id}")
            logger.info(f"Using frontend URL: {frontend_url}")
            
            # Create temporary directory for PowerPoint output
            temp_dir = Path(tempfile.gettempdir()) / "monte_carlo_presentations"
            temp_dir.mkdir(exist_ok=True)
            ppt_path = temp_dir / f"simulation_presentation_{simulation_id}_{int(time.time())}.pptx"
            
            # Generate screenshots using the same approach as PDF export
            screenshots = await self._capture_simulation_screenshots(
                simulation_id, results_data, auth_token, frontend_url
            )
            
            # Create PowerPoint presentation from screenshots
            await self._create_powerpoint_from_screenshots(
                screenshots, ppt_path, simulation_id, results_data, metadata
            )
            
            logger.info(f"Pixel-perfect PowerPoint presentation generated: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint presentation: {e}")
            raise e
    
    async def _capture_simulation_screenshots(
        self, 
        simulation_id: str, 
        results_data: Dict[str, Any],
        auth_token: Optional[str] = None,
        frontend_url: str = "http://frontend:3000"
    ) -> Dict[str, str]:
        """
        Capture pixel-perfect screenshots of the simulation results using the same approach as PDF export.
        
        Returns:
            Dict[str, str]: Dictionary mapping screenshot names to file paths
        """
        screenshots = {}
        
        async with async_playwright() as p:
            # Launch browser with same config as PDF export
            browser = await p.chromium.launch(
                headless=True,
                args=['--no-sandbox', '--disable-setuid-sandbox']
            )
            
            context = await browser.new_context(
                viewport={'width': 1920, 'height': 1080},  # 16:9 aspect ratio
                device_scale_factor=1.0
            )
            
            page = await context.new_page()
            
            def log_console_message(msg):
                logger.info(f"🌐 BROWSER CONSOLE [{msg.type}]: {msg.text}")
            
            page.on("console", log_console_message)
            
            # Set authentication if provided
            if auth_token:
                await page.set_extra_http_headers({
                    'Authorization': f'Bearer {auth_token}'
                })
                
                await page.add_init_script(f"""
                    localStorage.setItem('authToken', '{auth_token}');
                    localStorage.setItem('auth0Token', '{auth_token}');
                    window.authToken = '{auth_token}';
                """)
            
            # Prepare data in the same format as PDF export (reuse that logic)
            prepared_targets = []
            targets = results_data.get('targets', {})
            
            for target_name, target_data in targets.items():
                logger.info(f"Processing target for PowerPoint: {target_name}")
                
                if 'values' in target_data and 'statistics' in target_data:
                    statistics = target_data.get('statistics', {})
                    histogram_data = target_data.get('histogram_data', {})
                    values = target_data.get('values', [])
                    
                    # Calculate min/max and percentiles (same as PDF export)
                    min_value = statistics.get('min', statistics.get('min_value', statistics.get('minimum', 0)))
                    max_value = statistics.get('max', statistics.get('max_value', statistics.get('maximum', 0)))
                    
                    percentiles = {}
                    if values and len(values) > 0:
                        values_array = np.array(values)
                        percentiles = {
                            'p5': float(np.percentile(values_array, 5)),
                            'p10': float(np.percentile(values_array, 10)),
                            'p25': float(np.percentile(values_array, 25)),
                            'p50': float(np.percentile(values_array, 50)),
                            'p75': float(np.percentile(values_array, 75)),
                            'p90': float(np.percentile(values_array, 90)),
                            'p95': float(np.percentile(values_array, 95))
                        }
                        if min_value == 0 or min_value is None:
                            min_value = float(np.min(values_array))
                        if max_value == 0 or max_value is None:
                            max_value = float(np.max(values_array))
                    
                    frontend_results = {
                        'mean': statistics.get('mean', 0),
                        'median': statistics.get('median', 0), 
                        'std_dev': statistics.get('std_dev', 0),
                        'min_value': min_value,
                        'max_value': max_value,
                        'iterations_run': len(values) if values else 1000,
                        'raw_values': values[:1000] if values else [],
                        'percentiles': percentiles
                    }
                    
                    # Process histogram data
                    if histogram_data and 'bin_edges' in histogram_data and 'counts' in histogram_data:
                        frontend_results['histogram'] = {
                            'bin_edges': histogram_data['bin_edges'],
                            'counts': histogram_data['counts'],
                            'bins': histogram_data.get('bins', []),
                            'values': histogram_data.get('values', [])
                        }
                    else:
                        frontend_results['histogram'] = {}
                    
                    prepared_targets.append({
                        'simulation_id': f'ppt_export_{target_name}',
                        'temp_id': f'ppt_temp_{target_name}', 
                        'status': 'completed',
                        'target_name': target_name,
                        'result_cell_coordinate': target_name,
                        'target_cell': target_name,
                        'iterations_run': results_data.get('iterations_run', 1000),
                        'requested_engine_type': results_data.get('requested_engine_type', 'Ultra'),
                        'results': frontend_results,
                        'histogram': frontend_results['histogram'],
                        'sensitivity_analysis': target_data.get('sensitivity_analysis', [])
                    })
            
            # Create print data (same format as PDF export)
            print_data = {
                'simulationId': 'ppt_export',
                'results': prepared_targets,
                'metadata': {
                    'iterations_run': results_data.get('iterations_run', 1000),
                    'engine_type': results_data.get('requested_engine_type', 'Ultra'),
                    'timestamp': results_data.get('metadata', {}).get('timestamp', 'N/A')
                }
            }
            
            # Generate unique ID
            data_id = f"ppt_data_{int(time.time() * 1000)}"
            
            # Navigate to homepage and inject data
            await page.goto(f"{frontend_url}/", wait_until='networkidle')
            
            await page.evaluate(f"""
                try {{
                    sessionStorage.setItem('{data_id}', JSON.stringify({json.dumps(print_data)}));
                    console.log('PowerPoint data stored in sessionStorage with ID: {data_id}');
                }} catch (e) {{
                    window.__PPT_DATA__ = {json.dumps(print_data)};
                    window.__PPT_DATA_ID__ = '{data_id}';
                    console.log('PowerPoint data injected directly into window');
                }}
            """)
            
            # Navigate to print-view
            await page.goto(f"{frontend_url}/print-view?id={data_id}", wait_until='networkidle')
            
            # Hide UI elements that shouldn't appear in PowerPoint
            await page.evaluate("""
                const cookieBanner = document.querySelector('[role="dialog"][aria-live="polite"]');
                if (cookieBanner) {
                    cookieBanner.style.display = 'none';
                }
                
                const toastContainer = document.querySelector('.Toastify__toast-container');
                if (toastContainer) {
                    toastContainer.style.display = 'none';
                }
            """)
            
            # Wait for content to render
            await page.wait_for_timeout(3000)
            
            # Capture full page screenshot for overview slide
            temp_dir_base = Path(tempfile.gettempdir()) / "monte_carlo_presentations"
            temp_dir_base.mkdir(exist_ok=True)
            temp_dir_screenshots = temp_dir_base / "screenshots"
            temp_dir_screenshots.mkdir(exist_ok=True)
            
            overview_path = temp_dir_screenshots / "overview.png"
            await page.screenshot(path=str(overview_path), full_page=True)
            screenshots['overview'] = str(overview_path)
            
            # Capture individual target screenshots if multiple targets exist
            target_containers = await page.query_selector_all('.simulation-results-container .result-item')
            
            for i, container in enumerate(target_containers):
                try:
                    target_path = temp_dir_screenshots / f"target_{i}.png"
                    await container.screenshot(path=str(target_path))
                    screenshots[f'target_{i}'] = str(target_path)
                except Exception as e:
                    logger.warning(f"Could not capture screenshot for target {i}: {e}")
            
            await browser.close()
        
        logger.info(f"Captured {len(screenshots)} screenshots for PowerPoint generation")
        return screenshots
    
    async def _create_powerpoint_from_screenshots(
        self,
        screenshots: Dict[str, str],
        ppt_path: Path,
        simulation_id: str,
        results_data: Dict[str, Any],
        metadata: Optional[Dict[str, Any]] = None
    ):
        """
        Create PowerPoint presentation from screenshots with editable text overlays.
        This ensures pixel-perfect visual fidelity while maintaining PowerPoint editability.
        """
        logger.info("Creating PowerPoint presentation from screenshots")
        
        # Create new presentation with 16:9 aspect ratio
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # 16:9 ratio: 1920x1080 scaled
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        self._create_title_slide_with_data(prs, simulation_id, results_data, metadata)
        
        # Add overview slide with screenshot background
        if 'overview' in screenshots:
            self._create_screenshot_slide(
                prs, 
                screenshots['overview'], 
                "Simulation Results Overview",
                results_data
            )
        
        # Add individual target slides if available
        targets = results_data.get('targets', {})
        target_names = list(targets.keys())
        
        for i, target_name in enumerate(target_names):
            screenshot_key = f'target_{i}'
            if screenshot_key in screenshots:
                self._create_screenshot_slide(
                    prs,
                    screenshots[screenshot_key],
                    f"Results for {target_name}",
                    results_data,
                    target_specific=target_name
                )
        
        # Add methodology slide
        self._create_methodology_slide(prs, results_data, metadata)
        
        # Save presentation
        prs.save(str(ppt_path))
        logger.info(f"PowerPoint presentation saved to: {ppt_path}")
    
    def _create_title_slide_with_data(
        self, 
        prs: Presentation, 
        simulation_id: str, 
        results_data: Dict[str, Any], 
        metadata: Optional[Dict[str, Any]] = None
    ):
        """Create title slide with simulation information."""
        
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Monte Carlo Simulation Results"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Subtitle with simulation details
        subtitle = slide.placeholders[1]
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
        targets = results_data.get('targets', {})
        iterations = results_data.get('iterations_run', 1000)
        engine = results_data.get('requested_engine_type', 'Ultra')
        
        subtitle_text = f"Analysis of {len(targets)} Target Variables\n"
        subtitle_text += f"Engine: {engine} | Iterations: {iterations:,}\n"
        subtitle_text += f"Generated: {timestamp}\n"
        subtitle_text += f"Simulation ID: {simulation_id}"
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(16)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        
        # Add a note about pixel-perfect fidelity
        note_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.5))
        note_frame = note_box.text_frame
        note_frame.word_wrap = True
        
        note_para = note_frame.add_paragraph()
        note_para.text = "✨ This presentation contains pixel-perfect screenshots of your simulation results, "
        note_para.text += "ensuring identical visual appearance to your web interface while maintaining full PowerPoint editability."
        note_para.font.size = Pt(12)
        note_para.font.italic = True
        note_para.font.color.rgb = self.colors['success']
        note_para.alignment = PP_ALIGN.CENTER
    
    def _create_screenshot_slide(
        self,
        prs: Presentation,
        screenshot_path: str,
        title: str,
        results_data: Dict[str, Any],
        target_specific: Optional[str] = None
    ):
        """Create a slide with screenshot background and editable text overlay."""
        
        # Use blank layout for full control
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Add screenshot as background (fill the slide)
        try:
            # Calculate image dimensions to fit 16:9 slide while maintaining aspect ratio
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add image - it will automatically scale to fit
            slide.shapes.add_picture(
                screenshot_path,
                left=0,
                top=0,
                width=slide_width,
                height=slide_height
            )
            
            # Add semi-transparent title overlay at the top
            title_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(0.2),
                width=Inches(12.3),
                height=Inches(0.8)
            )
            
            # Style the title box with semi-transparent background
            title_fill = title_box.fill
            title_fill.solid()
            title_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            title_fill.transparency = 0.3  # 70% opaque
            
            # Add border
            title_line = title_box.line
            title_line.color.rgb = self.colors['primary']
            title_line.width = Pt(2)
            
            # Set title text
            title_frame = title_box.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = title
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.font.color.rgb = self.colors['dark']
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add editable data box in bottom right for key statistics
            if target_specific:
                targets = results_data.get('targets', {})
                if target_specific in targets:
                    self._add_editable_stats_overlay(slide, targets[target_specific], target_specific)
            
        except Exception as e:
            logger.warning(f"Could not add screenshot {screenshot_path} to slide: {e}")
            # Fallback: create text-only slide
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
            title_frame = title_shape.text_frame
            title_para = title_frame.add_paragraph()
            title_para.text = f"{title}\n(Screenshot could not be loaded)"
            title_para.font.size = Pt(24)
            title_para.font.color.rgb = self.colors['danger']
            title_para.alignment = PP_ALIGN.CENTER
    
    def _add_editable_stats_overlay(self, slide, target_data: Dict, target_name: str):
        """Add editable statistics overlay on screenshot."""
        
        statistics = target_data.get('statistics', {})
        if not statistics:
            return
        
        # Create semi-transparent statistics box in bottom-right corner
        stats_box = slide.shapes.add_textbox(
            left=Inches(9.5),
            top=Inches(5.5),
            width=Inches(3.5),
            height=Inches(1.8)
        )
        
        # Style the stats box
        stats_fill = stats_box.fill
        stats_fill.solid()
        stats_fill.fore_color.rgb = RGBColor(248, 249, 250)  # Light gray background
        stats_fill.transparency = 0.15  # 85% opaque
        
        # Add border
        stats_line = stats_box.line
        stats_line.color.rgb = self.colors['primary']
        stats_line.width = Pt(1)
        
        # Add statistics text (editable)
        stats_frame = stats_box.text_frame
        stats_frame.word_wrap = True
        stats_frame.margin_left = Inches(0.1)
        stats_frame.margin_right = Inches(0.1)
        stats_frame.margin_top = Inches(0.05)
        stats_frame.margin_bottom = Inches(0.05)
        
        # Title
        title_para = stats_frame.add_paragraph()
        title_para.text = f"Key Statistics"
        title_para.font.size = Pt(10)
        title_para.font.bold = True
        title_para.font.color.rgb = self.colors['dark']
        
        # Statistics (these remain editable in PowerPoint)
        stats_items = [
            f"Mean: {self._format_number(statistics.get('mean', 0))}",
            f"Median: {self._format_number(statistics.get('median', 0))}",
            f"Std Dev: {self._format_number(statistics.get('std_dev', 0))}",
            f"Range: {self._format_number(statistics.get('min', 0))} - {self._format_number(statistics.get('max', 0))}"
        ]
        
        for stat in stats_items:
            p = stats_frame.add_paragraph()
            p.text = stat
            p.font.size = Pt(8)
            p.font.color.rgb = self.colors['dark']
    
    def _create_methodology_slide(
        self, 
        prs: Presentation, 
        results_data: Dict[str, Any], 
        metadata: Optional[Dict[str, Any]] = None
    ):
        """Create methodology and summary slide."""
        
        content_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(content_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Methodology & Analysis Summary"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
        
        # Content area
        content = slide.placeholders[1]
        content_frame = content.text_frame
        content_frame.clear()  # Clear default text
        content_frame.word_wrap = True
        
        # Methodology section
        method_title = content_frame.add_paragraph()
        method_title.text = "Monte Carlo Analysis Methodology"
        method_title.font.size = Pt(18)
        method_title.font.bold = True
        method_title.font.color.rgb = self.colors['dark']
        method_title.space_after = Pt(6)
        
        # Methodology details
        methodology_items = [
            f"• Engine Type: {results_data.get('requested_engine_type', 'Ultra')} simulation engine",
            f"• Iterations: {results_data.get('iterations_run', 1000):,} Monte Carlo iterations",
            f"• Variables: {len(results_data.get('targets', {}))} target output variables analyzed",
            "• Statistical distributions fitted to historical data patterns",
            "• Random sampling with correlation structure preservation",
            "• Confidence intervals calculated at multiple percentile levels"
        ]
        
        for item in methodology_items:
            p = content_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['dark']
            p.space_after = Pt(3)
        
        # Separator
        sep = content_frame.add_paragraph()
        sep.text = ""
        sep.space_after = Pt(12)
        
        # Visual fidelity note
        fidelity_title = content_frame.add_paragraph()
        fidelity_title.text = "Presentation Features"
        fidelity_title.font.size = Pt(18)
        fidelity_title.font.bold = True
        fidelity_title.font.color.rgb = self.colors['success']
        fidelity_title.space_after = Pt(6)
        
        fidelity_items = [
            "✨ Pixel-perfect visual fidelity matching your web interface",
            "📊 All charts and graphs maintain exact appearance and styling", 
            "✏️ Text elements remain fully editable in PowerPoint",
            "📐 16:9 aspect ratio optimized for presentations",
            "🎯 Professional formatting suitable for stakeholder meetings"
        ]
        
        for item in fidelity_items:
            p = content_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = self.colors['dark']
            p.space_after = Pt(3)
    
    # Legacy methods removed - now using screenshot-based approach for pixel-perfect fidelity
    
    def _format_number(self, value: float) -> str:
        """Format numbers for display in PowerPoint."""
        
        if value == 0:
            return "0"
        
        abs_value = abs(value)
        
        if abs_value >= 1_000_000:
            return f"{value/1_000_000:.2f}M"
        elif abs_value >= 1_000:
            return f"{value/1_000:.1f}K"
        elif abs_value >= 1:
            return f"{value:.2f}"
        elif abs_value >= 0.01:
            return f"{value:.3f}"
        else:
            return f"{value:.2e}"

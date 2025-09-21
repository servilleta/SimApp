import os
import logging
import json
import tempfile
import time
from typing import Dict, Any, Optional, List
from datetime import datetime
from pathlib import Path
import io
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
import numpy as np

logger = logging.getLogger(__name__)

class PowerPointExportService:
    """Service for generating editable PowerPoint presentations from simulation results."""
    
    def __init__(self):
        self.logger = logger
        
        # Color scheme for consistency
        self.colors = {
            'primary': RGBColor(33, 150, 243),     # Blue
            'secondary': RGBColor(255, 152, 0),    # Orange  
            'success': RGBColor(76, 175, 80),      # Green
            'warning': RGBColor(255, 193, 7),      # Yellow
            'danger': RGBColor(244, 67, 54),       # Red
            'dark': RGBColor(33, 37, 41),          # Dark gray
            'light': RGBColor(248, 249, 250)       # Light gray
        }
    
    async def generate_powerpoint_presentation(
        self, 
        simulation_id: str, 
        results_data: Dict[str, Any],
        metadata: Optional[Dict[str, Any]] = None,
        auth_token: Optional[str] = None,
        frontend_url: str = "http://frontend:3000"
    ) -> str:
        """
        Generate hybrid PowerPoint presentation with both screenshot and editable versions.
        Provides exact visual fidelity AND fully editable elements in separate slides.
        
        Args:
            simulation_id: Unique identifier for the simulation
            results_data: Simulation results including targets and statistics
            metadata: Additional metadata about the simulation
            auth_token: Authentication token for frontend access
            frontend_url: URL of the frontend service
            
        Returns:
            str: Path to the generated PowerPoint file
        """
        try:
            logger.info(f"Starting hybrid PowerPoint generation for simulation {simulation_id}")
            
            # Create temporary directory for PowerPoint output
            temp_dir = Path(tempfile.gettempdir()) / "monte_carlo_presentations"
            temp_dir.mkdir(exist_ok=True)
            ppt_path = temp_dir / f"simulation_presentation_{simulation_id}_{int(time.time())}.pptx"
            
            # Create hybrid PowerPoint with both screenshot and editable versions
            await self._create_hybrid_powerpoint(
                ppt_path, simulation_id, results_data, metadata, auth_token, frontend_url
            )
            
            logger.info(f"Hybrid PowerPoint presentation generated: {ppt_path}")
            return str(ppt_path)
            
        except Exception as e:
            logger.error(f"Error generating PowerPoint presentation: {e}")
            raise e
    
    def _create_powerpoint_with_native_elements(
        self, 
        ppt_path: Path, 
        simulation_id: str, 
        results_data: Dict[str, Any], 
        metadata: Optional[Dict[str, Any]]
    ):
        """Create PowerPoint with fully editable native elements."""
        try:
            # Create new presentation with 16:9 aspect ratio
            prs = Presentation()
            prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
            prs.slide_height = Inches(7.5)
            
            # Add slides with native elements
            self._create_title_slide(prs, simulation_id, metadata)
            self._create_overview_slide(prs, results_data)
            
            # Create individual slides for each target
            targets = results_data.get('targets', {})
            for target_name, target_data in targets.items():
                self._create_target_slide(prs, target_name, target_data)
            
            # Add methodology slide
            self._create_methodology_slide(prs)
            
            # Save presentation
            prs.save(str(ppt_path))
            logger.info(f"PowerPoint saved with {len(prs.slides)} slides")
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint with native elements: {e}")
            raise e
    
    def _create_title_slide(self, prs: Presentation, simulation_id: str, metadata: Optional[Dict[str, Any]]):
        """Create title slide with simulation information."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Monte Carlo Simulation Results"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        
        # Subtitle with simulation details
        subtitle = slide.placeholders[1]
        subtitle_text = f"Simulation ID: {simulation_id}\n"
        
        if metadata:
            subtitle_text += f"Engine: {metadata.get('engine_type', 'Ultra')}\n"
            subtitle_text += f"Iterations: {metadata.get('iterations_run', 'N/A'):,}\n"
            subtitle_text += f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['primary']
    
    def _create_overview_slide(self, prs: Presentation, results_data: Dict[str, Any]):
        """Create overview slide with summary statistics."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Simulation Results Overview"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        
        # Create summary table
        targets = results_data.get('targets', {})
        if targets:
            self._add_summary_table(slide, targets)
        
        # Add box plot chart for comparison
        if len(targets) > 1:
            self._add_comparison_chart(slide, targets)
    
    def _add_summary_table(self, slide, targets: Dict[str, Any]):
        """Add summary statistics table."""
        # Table dimensions
        rows = len(targets) + 1  # +1 for header
        cols = 6  # Target, Mean, Median, Std Dev, Min, Max
        
        # Position table on left side
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(0.8 * rows)
        
        # Create table
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Header row
        headers = ['Target', 'Mean', 'Median', 'Std Dev', 'Min', 'Max']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.colors['primary']
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Data rows
        for row_idx, (target_name, target_data) in enumerate(targets.items(), 1):
            statistics = target_data.get('statistics', {})
            
            # Format values
            values = [
                target_name,
                f"{statistics.get('mean', 0):,.2f}",
                f"{statistics.get('median', 0):,.2f}",
                f"{statistics.get('std_dev', 0):,.2f}",
                f"{statistics.get('min_value', statistics.get('min', 0)):,.2f}",
                f"{statistics.get('max_value', statistics.get('max', 0)):,.2f}"
            ]
            
            for col_idx, value in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                
                # Alternate row colors
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.colors['light']
    
    def _add_comparison_chart(self, slide, targets: Dict[str, Any]):
        """Add comparison chart for multiple targets."""
        # Chart data
        chart_data = CategoryChartData()
        target_names = list(targets.keys())
        chart_data.categories = target_names
        
        # Add mean values series
        means = []
        for target_data in targets.values():
            statistics = target_data.get('statistics', {})
            means.append(statistics.get('mean', 0))
        
        chart_data.add_series('Mean Values', means)
        
        # Position chart on right side
        left = Inches(7)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        
        # Create chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
        )
        
        chart = chart_shape.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        
        # Customize chart
        chart.chart_title.text_frame.text = "Target Comparison"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    def _create_target_slide(self, prs: Presentation, target_name: str, target_data: Dict[str, Any]):
        """Create detailed slide for a specific target."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = f"Results for {target_name}"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        
        # Statistics box
        self._add_statistics_textbox(slide, target_data)
        
        # Histogram chart
        histogram_data = target_data.get('histogram_data', {})
        if histogram_data and 'bin_edges' in histogram_data and 'counts' in histogram_data:
            self._add_histogram_chart(slide, histogram_data, target_name)
        
        # Sensitivity analysis
        sensitivity_data = target_data.get('sensitivity_analysis', [])
        if sensitivity_data:
            self._add_sensitivity_chart(slide, sensitivity_data)
    
    def _add_statistics_textbox(self, slide, target_data: Dict[str, Any]):
        """Add statistics text box."""
        statistics = target_data.get('statistics', {})
        
        # Position on left side
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(3.5)
        height = Inches(3)
        
        # Create text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Add statistics
        stats_text = "Key Statistics:\n\n"
        stats_text += f"Mean: {statistics.get('mean', 0):,.2f}\n"
        stats_text += f"Median: {statistics.get('median', 0):,.2f}\n"
        stats_text += f"Std Dev: {statistics.get('std_dev', 0):,.2f}\n"
        stats_text += f"Min: {statistics.get('min_value', statistics.get('min', 0)):,.2f}\n"
        stats_text += f"Max: {statistics.get('max_value', statistics.get('max', 0)):,.2f}\n"
        
        # Add percentiles if available
        percentiles = statistics.get('percentiles', {})
        if percentiles:
            stats_text += f"\nPercentiles:\n"
            stats_text += f"P5: {percentiles.get('p5', 0):,.2f}\n"
            stats_text += f"P25: {percentiles.get('p25', 0):,.2f}\n"
            stats_text += f"P75: {percentiles.get('p75', 0):,.2f}\n"
            stats_text += f"P95: {percentiles.get('p95', 0):,.2f}\n"
        
        text_frame.text = stats_text
        
        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            if paragraph.text.startswith("Key Statistics:"):
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.color.rgb = self.colors['primary']
    
    def _add_histogram_chart(self, slide, histogram_data: Dict[str, Any], target_name: str):
        """Add histogram chart."""
        bin_edges = histogram_data.get('bin_edges', [])
        counts = histogram_data.get('counts', [])
        
        if not bin_edges or not counts:
            return
        
        # Create chart data
        chart_data = CategoryChartData()
        
        # Create bin labels (midpoints)
        bin_labels = []
        for i in range(len(bin_edges) - 1):
            midpoint = (bin_edges[i] + bin_edges[i + 1]) / 2
            bin_labels.append(f"{midpoint:.1f}")
        
        chart_data.categories = bin_labels
        chart_data.add_series('Frequency', counts)
        
        # Position chart
        left = Inches(4.5)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(3)
        
        # Create chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
        )
        
        chart = chart_shape.chart
        chart.has_legend = False
        
        # Customize chart
        chart.chart_title.text_frame.text = f"Distribution of {target_name}"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    def _add_sensitivity_chart(self, slide, sensitivity_data: List[Dict[str, Any]]):
        """Add sensitivity analysis chart."""
        if not sensitivity_data:
            return
        
        # Create chart data
        chart_data = CategoryChartData()
        
        variables = [item.get('variable', 'Unknown') for item in sensitivity_data]
        impacts = [abs(item.get('impact', 0)) for item in sensitivity_data]
        
        chart_data.categories = variables
        chart_data.add_series('Impact', impacts)
        
        # Position chart
        left = Inches(4.5)
        top = Inches(4.8)
        width = Inches(8)
        height = Inches(2.5)
        
        # Create chart
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
        )
        
        chart = chart_shape.chart
        chart.has_legend = False
        
        # Customize chart
        chart.chart_title.text_frame.text = "Variable Impact Analysis"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
    
    def _create_methodology_slide(self, prs: Presentation):
        """Create methodology slide."""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title = slide.shapes.title
        title.text = "Methodology & Features"
        title.text_frame.paragraphs[0].font.size = Pt(36)
        title.text_frame.paragraphs[0].font.color.rgb = self.colors['dark']
        
        # Content
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        methodology_text = """Monte Carlo Simulation Methodology:

• Ultra Engine: High-performance GPU-accelerated simulation engine
• Statistical Analysis: Comprehensive statistical measures including percentiles
• Sensitivity Analysis: Variable impact assessment using correlation analysis
• Visualization: Interactive charts and histograms for result interpretation

PowerPoint Features:
• Fully Editable Elements: All charts, tables, and text can be modified
• Native PowerPoint Charts: Use PowerPoint's chart editing capabilities
• Professional Layout: 16:9 aspect ratio optimized for presentations
• Data Integration: Statistics and results embedded as editable elements

This presentation contains live data that can be updated and customized using PowerPoint's built-in editing tools."""
        
        text_frame.text = methodology_text
        
        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.alignment = PP_ALIGN.LEFT
            if paragraph.text.endswith(":"):
                paragraph.font.bold = True
                paragraph.font.color.rgb = self.colors['primary']
